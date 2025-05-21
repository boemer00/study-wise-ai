from datetime import datetime
from typing import List, Optional, Dict, Any
import os
import mimetypes
from pathlib import Path
import requests
from dataclasses import dataclass
import fitz  # PyMuPDF
import docx
from langchain.text_splitter import RecursiveCharacterTextSplitter
import urllib.parse
import tempfile
import shutil


@dataclass
class UploadedDocument:
    """Represents an uploaded document with metadata."""
    file_name: str
    file_type: str
    upload_date: datetime
    content: str
    chunks: List[str] = None


class DocumentUploader:
    """Handles document uploads from various sources and formats."""

    def __init__(self, max_file_size_mb: int = 10):
        """
        Initialize the document uploader.

        Args:
            max_file_size_mb: Maximum allowed file size in megabytes
        """
        self.max_file_size_mb = max_file_size_mb
        self._supported_formats = {
            "pdf": "application/pdf",
            "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "doc": "application/msword",
            "txt": "text/plain",
            "md": "text/markdown",
            "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "ppt": "application/vnd.ms-powerpoint"
        }
        self.processor = DocumentProcessor()

    def upload_from_file(self, file_path: str) -> UploadedDocument:
        """
        Upload and extract text from a local file.

        Args:
            file_path: Path to the local file

        Returns:
            UploadedDocument with extracted content

        Raises:
            ValueError: If file format is not supported or file size exceeds limit
            FileNotFoundError: If file does not exist
        """
        file_path = Path(file_path)

        # Check if file exists
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        # Check file size
        file_size_mb = file_path.stat().st_size / (1024 * 1024)
        if file_size_mb > self.max_file_size_mb:
            raise ValueError(f"File size ({file_size_mb:.2f} MB) exceeds the maximum allowed size ({self.max_file_size_mb} MB)")

        # Check file format
        file_extension = file_path.suffix.lower().lstrip('.')
        if file_extension not in self._supported_formats:
            raise ValueError(f"Unsupported file format: {file_extension}. Supported formats: {', '.join(self._supported_formats.keys())}")

        # Process the document to extract text content
        content = self.processor.process_document(str(file_path))

        # Create and return the uploaded document
        return UploadedDocument(
            file_name=file_path.name,
            file_type=file_extension,
            upload_date=datetime.now(),
            content=content,
            chunks=self.processor.chunk_text(content)
        )

    def upload_from_url(self, url: str, download_dir: Optional[str] = None) -> UploadedDocument:
        """
        Download, upload and extract text from a document URL.

        Args:
            url: URL to the document
            download_dir: Directory to save the downloaded file (optional)

        Returns:
            UploadedDocument with extracted content

        Raises:
            ValueError: If URL format is invalid, file format is not supported,
                        or file size exceeds limit
            ConnectionError: If download fails
        """
        # Validate URL
        try:
            parsed_url = urllib.parse.urlparse(url)
            if not parsed_url.scheme or not parsed_url.netloc:
                raise ValueError(f"Invalid URL: {url}")
        except Exception as e:
            raise ValueError(f"Invalid URL: {url}. Error: {str(e)}")

        # Get the filename from the URL
        file_name = os.path.basename(parsed_url.path)
        if not file_name:
            file_name = "downloaded_document"

        # Create a temporary directory if download_dir is not provided
        temp_dir = None
        if not download_dir:
            temp_dir = tempfile.mkdtemp()
            download_dir = temp_dir
        else:
            os.makedirs(download_dir, exist_ok=True)

        download_path = os.path.join(download_dir, file_name)

        try:
            # Download the file with a streaming request to check size during download
            response = requests.get(url, stream=True)
            if response.status_code != 200:
                raise ConnectionError(f"Failed to download file from {url}. Status code: {response.status_code}")

            # Check content type
            content_type = response.headers.get('content-type', '').lower()
            file_extension = mimetypes.guess_extension(content_type)

            if not file_extension:
                # Try to get extension from URL if content-type doesn't help
                file_extension = os.path.splitext(file_name)[1].lower()

            if file_extension.startswith('.'):
                file_extension = file_extension[1:]

            if file_extension not in self._supported_formats:
                raise ValueError(f"Unsupported file format: {file_extension or 'unknown'}. Supported formats: {', '.join(self._supported_formats.keys())}")

            # Check file size from headers
            content_length = response.headers.get('content-length')
            if content_length and int(content_length) / (1024 * 1024) > self.max_file_size_mb:
                raise ValueError(f"File size exceeds the maximum allowed size ({self.max_file_size_mb} MB)")

            # Download the file
            with open(download_path, 'wb') as f:
                downloaded_size = 0
                for chunk in response.iter_content(chunk_size=8192):
                    downloaded_size += len(chunk)
                    if downloaded_size / (1024 * 1024) > self.max_file_size_mb:
                        raise ValueError(f"File size exceeds the maximum allowed size ({self.max_file_size_mb} MB)")
                    f.write(chunk)

            # Process the downloaded file
            result = self.upload_from_file(download_path)

            return result

        except Exception as e:
            raise e
        finally:
            # Clean up temporary directory if it was created
            if temp_dir:
                shutil.rmtree(temp_dir, ignore_errors=True)

    def upload_from_text(self, text: str, file_name: str = "pasted_text.txt") -> UploadedDocument:
        """
        Accept directly pasted text.

        Args:
            text: The text content
            file_name: Name to assign to the document

        Returns:
            UploadedDocument with the provided content
        """
        if not text:
            raise ValueError("Text content cannot be empty")

        # Create uploaded document directly from the provided text
        document = UploadedDocument(
            file_name=file_name,
            file_type="txt",
            upload_date=datetime.now(),
            content=text,
            chunks=self.processor.chunk_text(text)
        )

        return document

    def get_supported_formats(self) -> List[str]:
        """
        Return list of supported file formats.

        Returns:
            List of supported extensions
        """
        return list(self._supported_formats.keys())


class DocumentProcessor:
    """Processes uploaded documents and extracts text content."""

    def __init__(self):
        """Initialize the document processor."""
        self._extractors = {
            "pdf": PDFExtractor(),
            "docx": DocxExtractor(),
            "doc": DocxExtractor(),  # Use the same extractor for .doc files
            "txt": TextExtractor(),
            "md": TextExtractor(),   # Use the same extractor for markdown files
            "pptx": PPTExtractor(),
            "ppt": PPTExtractor()    # Use the same extractor for .ppt files
        }

    def process_document(self, file_path: str) -> str:
        """
        Extract text from the document based on its format.

        Args:
            file_path: Path to the document

        Returns:
            Extracted text content as string

        Raises:
            ValueError: If file format is not supported
            FileNotFoundError: If file does not exist
        """
        file_path = Path(file_path)

        # Check if file exists
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        # Get file extension
        file_extension = file_path.suffix.lower().lstrip('.')

        # Check if the format is supported
        if file_extension not in self._extractors:
            raise ValueError(f"Unsupported document format: {file_extension}")

        # Extract text using the appropriate extractor
        extractor = self._extractors[file_extension]
        return extractor.extract(str(file_path))

    def chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """
        Split extracted text into manageable chunks for processing.

        Args:
            text: The text to split
            chunk_size: Maximum size of each chunk
            overlap: Overlap between consecutive chunks

        Returns:
            List of text chunks
        """
        if not text:
            return []

        # Use langchain's text splitter to split the text into chunks
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=overlap,
            separators=["\n\n", "\n", " ", ""]
        )

        chunks = splitter.split_text(text)
        return chunks


# Format-specific extractors

class PDFExtractor:
    """Extracts text from PDF documents."""

    def extract(self, file_path: str) -> str:
        """
        Extract text from a PDF file.

        Args:
            file_path: Path to the PDF file

        Returns:
            Extracted text
        """
        try:
            # Open the PDF file
            doc = fitz.open(file_path)

            # Extract text from each page and join with double newlines
            text = "\n\n".join(page.get_text("text") for page in doc)

            return text
        except Exception as e:
            raise ValueError(f"Failed to extract text from PDF: {str(e)}")


class DocxExtractor:
    """Extracts text from DOCX documents."""

    def extract(self, file_path: str) -> str:
        """
        Extract text from a DOCX file.

        Args:
            file_path: Path to the DOCX file

        Returns:
            Extracted text
        """
        try:
            # Load the DOCX document
            doc = docx.Document(file_path)

            # Extract text from paragraphs
            paragraphs = [p.text for p in doc.paragraphs]

            # Join paragraphs with double newlines
            text = "\n\n".join(paragraphs)

            return text
        except Exception as e:
            raise ValueError(f"Failed to extract text from DOCX: {str(e)}")


class TextExtractor:
    """Extracts text from plain text documents."""

    def extract(self, file_path: str) -> str:
        """
        Extract text from a plain text file.

        Args:
            file_path: Path to the text file

        Returns:
            Extracted text
        """
        try:
            # Open and read the text file
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()

            return text
        except UnicodeDecodeError:
            # Try with a different encoding if UTF-8 fails
            try:
                with open(file_path, 'r', encoding='latin-1') as f:
                    text = f.read()
                return text
            except Exception as e:
                raise ValueError(f"Failed to extract text from text file: {str(e)}")
        except Exception as e:
            raise ValueError(f"Failed to extract text from text file: {str(e)}")


class PPTExtractor:
    """Extracts text from PowerPoint presentations."""

    def extract(self, file_path: str) -> str:
        """
        Extract text from a PPT/PPTX file.

        Args:
            file_path: Path to the presentation file

        Returns:
            Extracted text
        """
        try:
            # For PPTX, use the python-pptx library
            from pptx import Presentation

            # Load the presentation
            presentation = Presentation(file_path)

            # Extract text from each slide
            text_parts = []

            for i, slide in enumerate(presentation.slides):
                slide_text = []
                slide_text.append(f"Slide {i+1}:")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)

                text_parts.append("\n".join(slide_text))

            # Join slide texts with double newlines
            text = "\n\n".join(text_parts)

            return text
        except ImportError:
            raise ValueError("python-pptx package is required to extract text from PowerPoint files. Please install it with 'pip install python-pptx'")
        except Exception as e:
            raise ValueError(f"Failed to extract text from PowerPoint file: {str(e)}")
